import os
import io
import json
import logging
import re
import traceback
import streamlit as st

logging.basicConfig(level=logging.INFO)

from ai_client import DEFAULT_MODEL, AIClientError, chat_completion
from data_loader import build_context, load_file

SYSTEM_PROMPT = """
You are an AI document analyst.

Answer questions only using the uploaded document context.

Rules:
- Search the uploaded document before answering.
- Give only information available in the document.
- Do not use external knowledge or assumptions.
- Do not recommend checking HR or other sources.
- If the document has relevant information under a different section name,
  use that section to answer the question.
- If the answer is partial, give the available details and say what specific
  part is not available.
- When the user asks to share, show, display, provide or explain a policy,
  subpolicy, procedure or form, present the main document information clearly.
- For policies and subpolicies, include the title, page or section, purpose,
  applicability, eligibility, rules, workflow and limits when present.
- For forms, include the form name, page or annexure, purpose, who initiates it,
  who approves it, required fields, process flow and any notes/instructions when
  present.
- If the answer is missing, say:
"The information is not available in the uploaded document."
- Mention page number or section name when available.
- Keep answers clear, direct and relevant to the user's question.
"""


FORM_GENERATION_PROMPT = """
You are an HR document generation assistant.

Use only the uploaded document context. Create a practical HR form for the
employee request. Return only valid JSON with this schema:
{
  "title": "Form title",
  "summary": "Short explanation of what policy information was used.",
  "eligibility": ["item"],
  "duration": ["item"],
  "approval_workflow": ["item"],
  "fields": ["field label"],
  "notes": ["item"]
}

Rules:
- Include fields that the policy requires or clearly implies.
- Include Employee Name, Employee ID and Department if the request is for an
  employee application form.
- Include approval fields when the policy mentions an approval workflow.
- Do not invent policy limits, eligibility rules or duration values.
- If a detail is not available, put "Not available in uploaded document."
"""


PRESENTATION_GENERATION_PROMPT = """
You are an HR presentation generation assistant.

Use only the uploaded document context. Create concise slide content for the
employee request. Return only valid JSON with this schema:
{
  "title": "Presentation title",
  "summary": "Short explanation of what policy information was used.",
  "slides": [
    {
      "title": "Slide title",
      "bullets": ["short bullet"],
      "speaker_notes": "optional note"
    }
  ],
  "notes": ["item"]
}

Rules:
- Make a logical slide-to-slide story.
- Keep each slide focused with short bullets.
- Cover eligibility, duration, process, approval workflow, employee
  responsibilities and important limits when present in the document.
- Do not invent policy limits, eligibility rules or duration values.
- If a detail is not available, put "Not available in uploaded document."
"""



# Vision-capable models that accept image_url content via the standard
# /v1/chat/completions endpoint. NOTE: nvidia/nemotron-ocr-v1 is deliberately
# excluded — it's a specialized NeMo Retriever OCR microservice with its own
# detector/recognizer API, not a chat-completions model, and calling it
# through chat_completion() will hang or fail rather than return text.
MODEL1_OPTIONS = [
    "meta/llama-3.2-11b-vision-instruct",
    "meta/llama-3.2-90b-vision-instruct",
    "nvidia/llama-3.1-nemotron-nano-vl-8b-v1",
]

# General-purpose text/instruct models for document generation (forms,
# policies, presentations). meta/llama-3.1-8b-instruct is the default: it's
# fast (responses in a few seconds) and is the model this app's NVIDIA key
# is provisioned for — see README.md.
MODEL2_OPTIONS = [
    DEFAULT_MODEL,
    "meta/llama-3.1-70b-instruct",
    "meta/llama-3.3-70b-instruct",
    "abacusai/dracarys-llama-3.1-70b-instruct",
]



def is_document_generation_request(question):
    text = question.lower()
    action_words = ("generate", "create", "prepare", "make", "draft")
    document_words = ("form", "application", "letter", "document", "presentation", "slides", "pdf", "docx", "pptx")
    return any(word in text for word in action_words) and any(word in text for word in document_words)


def is_presentation_request(question):
    text = question.lower()
    return any(word in text for word in ("presentation", "slides", "slide deck", "ppt", "pptx", "powerpoint"))


def filename_from_title(title):
    cleaned = re.sub(r"[^a-zA-Z0-9]+", "_", title.lower()).strip("_")
    return cleaned or "generated_document"


def extract_json_object(text):
    try:
        return json.loads(text)
    except json.JSONDecodeError:
        match = re.search(r"\{.*\}", text, flags=re.DOTALL)
        if not match:
            raise
        return json.loads(match.group(0))


def list_from_value(value):
    if value is None:
        return []
    if isinstance(value, list):
        return [str(item).strip() for item in value if str(item).strip()]
    text = str(value).strip()
    return [text] if text else []


def normalize_form_data(data, question):
    title = str(data.get("title") or question.strip().title())
    fields = data.get("fields") or []
    required_defaults = ["Employee Name", "Employee ID", "Department"]

    normalized_fields = []
    for field in [*required_defaults, *fields]:
        label = str(field).strip().rstrip(":")
        if label and label not in normalized_fields:
            normalized_fields.append(label)

    return {
        "title": title,
        "summary": str(data.get("summary") or ""),
        "eligibility": list_from_value(data.get("eligibility")),
        "duration": list_from_value(data.get("duration")),
        "approval_workflow": list_from_value(data.get("approval_workflow")),
        "fields": normalized_fields,
        "notes": list_from_value(data.get("notes")),
    }


def normalize_presentation_data(data, question):
    title = str(data.get("title") or question.strip().title())
    slides = []
    for index, slide in enumerate(data.get("slides") or [], start=1):
        if not isinstance(slide, dict):
            continue
        slide_title = str(slide.get("title") or f"Slide {index}").strip()
        bullets = list_from_value(slide.get("bullets"))
        speaker_notes = str(slide.get("speaker_notes") or "").strip()
        slides.append(
            {
                "title": slide_title,
                "bullets": bullets[:6],
                "speaker_notes": speaker_notes,
            }
        )

    if not slides:
        slides = [
            {
                "title": title,
                "bullets": ["Not available in uploaded document."],
                "speaker_notes": "",
            }
        ]

    return {
        "title": title,
        "summary": str(data.get("summary") or ""),
        "slides": slides[:12],
        "notes": list_from_value(data.get("notes")),
    }


def render_form_preview(form_data):
    lines = [f"### {form_data['title']}"]
    if form_data["summary"]:
        lines.extend(["", form_data["summary"]])

    for heading, key in (
        ("Eligibility", "eligibility"),
        ("Duration", "duration"),
        ("Approval Workflow", "approval_workflow"),
        ("Notes", "notes"),
    ):
        if form_data[key]:
            lines.extend(["", f"**{heading}**"])
            lines.extend(f"- {item}" for item in form_data[key])

    lines.extend(["", "**Application Form**"])
    lines.extend(f"{field}:" for field in form_data["fields"])
    return "\n".join(lines)


def render_presentation_preview(presentation_data):
    lines = [f"### {presentation_data['title']}"]
    if presentation_data["summary"]:
        lines.extend(["", presentation_data["summary"]])

    for index, slide in enumerate(presentation_data["slides"], start=1):
        lines.extend(["", f"**Slide {index}: {slide['title']}**"])
        lines.extend(f"- {bullet}" for bullet in slide["bullets"])

    if presentation_data["notes"]:
        lines.extend(["", "**Notes**"])
        lines.extend(f"- {item}" for item in presentation_data["notes"])
    return "\n".join(lines)


def create_docx_form(form_data):
    try:
        import docx
    except ImportError as exc:
        raise RuntimeError("DOCX generation needs python-docx. Run: pip install -r requirements.txt") from exc

    document = docx.Document()
    document.add_heading(form_data["title"], level=1)
    if form_data["summary"]:
        document.add_paragraph(form_data["summary"])

    for heading, key in (
        ("Eligibility", "eligibility"),
        ("Duration", "duration"),
        ("Approval Workflow", "approval_workflow"),
        ("Notes", "notes"),
    ):
        if form_data[key]:
            document.add_heading(heading, level=2)
            for item in form_data[key]:
                document.add_paragraph(item, style="List Bullet")

    document.add_heading("Application Form", level=2)
    table = document.add_table(rows=1, cols=2)
    table.style = "Table Grid"
    header_cells = table.rows[0].cells
    header_cells[0].text = "Field"
    header_cells[1].text = "Employee Input / Approval"

    for field in form_data["fields"]:
        cells = table.add_row().cells
        cells[0].text = field
        cells[1].text = ""

    buffer = io.BytesIO()
    document.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()


def create_docx_presentation(presentation_data):
    try:
        import docx
    except ImportError as exc:
        raise RuntimeError("DOCX generation needs python-docx. Run: pip install -r requirements.txt") from exc

    document = docx.Document()
    document.add_heading(presentation_data["title"], level=1)
    if presentation_data["summary"]:
        document.add_paragraph(presentation_data["summary"])

    for index, slide in enumerate(presentation_data["slides"], start=1):
        document.add_heading(f"Slide {index}: {slide['title']}", level=2)
        for bullet in slide["bullets"]:
            document.add_paragraph(bullet, style="List Bullet")
        if slide["speaker_notes"]:
            document.add_paragraph(f"Speaker notes: {slide['speaker_notes']}")

    if presentation_data["notes"]:
        document.add_heading("Notes", level=2)
        for item in presentation_data["notes"]:
            document.add_paragraph(item, style="List Bullet")

    buffer = io.BytesIO()
    document.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()


def create_pptx_presentation(presentation_data):
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt
    except ImportError as exc:
        raise RuntimeError("PPTX generation needs python-pptx. Run: pip install -r requirements.txt") from exc

    deck = Presentation()
    deck.slide_width = Inches(13.333)
    deck.slide_height = Inches(7.5)

    title_slide = deck.slides.add_slide(deck.slide_layouts[6])
    title_box = title_slide.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(11.8), Inches(1.2))
    title_frame = title_box.text_frame
    title_frame.clear()
    title_para = title_frame.paragraphs[0]
    title_run = title_para.add_run()
    title_run.text = presentation_data["title"]
    title_para.alignment = PP_ALIGN.CENTER
    title_run.font.size = Pt(44)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(29, 36, 53)

    if presentation_data["summary"]:
        summary_box = title_slide.shapes.add_textbox(Inches(1.4), Inches(3.55), Inches(10.6), Inches(1.1))
        summary_frame = summary_box.text_frame
        summary_frame.word_wrap = True
        summary_frame.clear()
        summary_para = summary_frame.paragraphs[0]
        summary_run = summary_para.add_run()
        summary_run.text = presentation_data["summary"]
        summary_para.alignment = PP_ALIGN.CENTER
        summary_run.font.size = Pt(20)
        summary_run.font.color.rgb = RGBColor(72, 82, 102)

    for slide_data in presentation_data["slides"]:
        slide = deck.slides.add_slide(deck.slide_layouts[6])
        title_box = slide.shapes.add_textbox(Inches(0.65), Inches(0.45), Inches(12.0), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.clear()
        title_run = title_frame.paragraphs[0].add_run()
        title_run.text = slide_data["title"]
        title_run.font.size = Pt(34)
        title_run.font.bold = True
        title_run.font.color.rgb = RGBColor(29, 36, 53)

        body_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.55), Inches(11.25), Inches(4.9))
        body_frame = body_box.text_frame
        body_frame.word_wrap = True
        body_frame.clear()

        for bullet_index, bullet in enumerate(slide_data["bullets"] or ["Not available in uploaded document."]):
            paragraph = body_frame.paragraphs[0] if bullet_index == 0 else body_frame.add_paragraph()
            paragraph.text = bullet
            paragraph.level = 0
            paragraph.font.size = Pt(22)
            paragraph.font.color.rgb = RGBColor(42, 50, 68)

        if slide_data["speaker_notes"]:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data["speaker_notes"]

    buffer = io.BytesIO()
    deck.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()


def create_pdf_form(form_data):
    try:
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.lib.units import inch
        from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle
        from xml.sax.saxutils import escape
    except ImportError as exc:
        raise RuntimeError("PDF generation needs reportlab. Run: pip install -r requirements.txt") from exc

    buffer = io.BytesIO()
    doc = SimpleDocTemplate(
        buffer,
        pagesize=A4,
        rightMargin=0.6 * inch,
        leftMargin=0.6 * inch,
        topMargin=0.6 * inch,
        bottomMargin=0.6 * inch,
    )
    styles = getSampleStyleSheet()
    story = [Paragraph(escape(form_data["title"]), styles["Title"])]

    if form_data["summary"]:
        story.extend([Spacer(1, 10), Paragraph(escape(form_data["summary"]), styles["BodyText"])])

    for heading, key in (
        ("Eligibility", "eligibility"),
        ("Duration", "duration"),
        ("Approval Workflow", "approval_workflow"),
        ("Notes", "notes"),
    ):
        if form_data[key]:
            story.extend([Spacer(1, 12), Paragraph(escape(heading), styles["Heading2"])])
            for item in form_data[key]:
                story.append(Paragraph(escape(f"- {item}"), styles["BodyText"]))

    story.extend([Spacer(1, 14), Paragraph("Application Form", styles["Heading2"])])
    table_data = [["Field", "Employee Input / Approval"]]
    table_data.extend([[field, ""] for field in form_data["fields"]])
    table = Table(table_data, colWidths=[2.4 * inch, 4.2 * inch], rowHeights=0.36 * inch)
    table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#e9edf5")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.HexColor("#1d2435")),
                ("GRID", (0, 0), (-1, -1), 0.7, colors.HexColor("#9aa4b2")),
                ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ]
        )
    )
    story.append(table)
    doc.build(story)
    buffer.seek(0)
    return buffer.getvalue()


def create_pdf_presentation(presentation_data):
    try:
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.lib.units import inch
        from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle
        from xml.sax.saxutils import escape
    except ImportError as exc:
        raise RuntimeError("PDF generation needs reportlab. Run: pip install -r requirements.txt") from exc

    buffer = io.BytesIO()
    doc = SimpleDocTemplate(
        buffer,
        pagesize=A4,
        rightMargin=0.6 * inch,
        leftMargin=0.6 * inch,
        topMargin=0.6 * inch,
        bottomMargin=0.6 * inch,
    )
    styles = getSampleStyleSheet()
    story = [Paragraph(escape(presentation_data["title"]), styles["Title"])]

    if presentation_data["summary"]:
        story.extend([Spacer(1, 10), Paragraph(escape(presentation_data["summary"]), styles["BodyText"])])

    for index, slide in enumerate(presentation_data["slides"], start=1):
        story.extend([Spacer(1, 12), Paragraph(escape(f"Slide {index}: {slide['title']}"), styles["Heading2"])])
        rows = [[Paragraph(escape(bullet), styles["BodyText"])] for bullet in slide["bullets"]]
        if not rows:
            rows = [[Paragraph("Not available in uploaded document.", styles["BodyText"])]]
        table = Table(rows, colWidths=[6.6 * inch])
        table.setStyle(
            TableStyle(
                [
                    ("BOX", (0, 0), (-1, -1), 0.5, colors.HexColor("#c8d0dc")),
                    ("INNERGRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#e1e5ec")),
                    ("VALIGN", (0, 0), (-1, -1), "TOP"),
                    ("LEFTPADDING", (0, 0), (-1, -1), 8),
                    ("RIGHTPADDING", (0, 0), (-1, -1), 8),
                    ("TOPPADDING", (0, 0), (-1, -1), 6),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
                ]
            )
        )
        story.append(table)
        if slide["speaker_notes"]:
            story.extend([Spacer(1, 6), Paragraph(escape(f"Speaker notes: {slide['speaker_notes']}"), styles["Italic"])])

    if presentation_data["notes"]:
        story.extend([Spacer(1, 12), Paragraph("Notes", styles["Heading2"])])
        for item in presentation_data["notes"]:
            story.append(Paragraph(escape(f"- {item}"), styles["BodyText"]))

    doc.build(story)
    buffer.seek(0)
    return buffer.getvalue()


st.set_page_config(
    page_title="AI Data Analyst - Universal Data Explorer",
    page_icon="🤖",
    layout="wide"
)


st.markdown("""
<style>

.block-container{
max-width:1500px;
padding-top:3rem;
padding-left:5rem;
padding-right:4rem;
}

[data-testid="stSidebar"]{
background:#f1f3f7;
}

.app-title{
font-size:48px;
font-weight:800;
color:#292b39;
}

.section-title{
font-size:32px;
font-weight:750;
color:#1d2435;
}

.chat-title{
font-size:38px;
font-weight:800;
color:#253047;
margin-top:45px;
}

[data-testid="stFileUploader"] section{
background:#f0f2f6;
border:none;
border-radius:10px;
min-height:120px;
}

.load-first{
background:#e7f1ff;
padding:18px;
border-radius:8px;
color:#07549b;
font-size:18px;
}

.loaded-box{
background:#eef8f6;
padding:15px;
border-radius:8px;
border:1px solid #cce5df;
}

</style>
""", unsafe_allow_html=True)


if "files" not in st.session_state:
    st.session_state.files=[]

if "file_keys" not in st.session_state:
    st.session_state.file_keys=set()

if "history" not in st.session_state:
    st.session_state.history=[]

if "ocr_model" not in st.session_state or st.session_state.ocr_model not in MODEL1_OPTIONS:
    st.session_state.ocr_model=MODEL1_OPTIONS[0]

if "llm_model" not in st.session_state or st.session_state.llm_model not in MODEL2_OPTIONS:
    st.session_state.llm_model=DEFAULT_MODEL



with st.sidebar:

    st.markdown("# ⚙️ Configuration")

    st.write("Models")

    st.session_state.ocr_model = st.selectbox(
        "OCR / vision model",
        MODEL1_OPTIONS,
        index=MODEL1_OPTIONS.index(st.session_state.ocr_model)
    )

    st.session_state.llm_model = st.selectbox(
        "LLM / chat model",
        MODEL2_OPTIONS,
        index=MODEL2_OPTIONS.index(st.session_state.llm_model)
    )

    ocr_api_key = st.text_input(
        "NVIDIA OCR API Key",
        type="password"
    )

    llm_api_key = st.text_input(
        "NVIDIA LLM API Key",
        type="password"
    )

    st.divider()

    st.markdown("## 📁 Data Source")

    st.radio(
        "Choose how to load data:",
        ["📤 Upload File"]
    )

    st.write("Upload CSV, Excel, PDF or any file")
    st.write("to analyze")

    if st.session_state.files:

        st.divider()

        st.write("### Loaded Files")

        for f in st.session_state.files:
            st.caption(f.filename)


    if st.button(
        "Clear chat and files",
        use_container_width=True
    ):
        st.session_state.files=[]
        st.session_state.file_keys=set()
        st.session_state.history=[]
        st.rerun()



st.markdown(
"<div class='app-title'>🤖 AI Data Analyst - Universal Data Explorer 🔗</div>",
unsafe_allow_html=True
)


st.markdown(
"<div class='section-title'>Intelligent Analysis for ANY Dataset | AI-Powered Insights</div>",
unsafe_allow_html=True
)


st.write("Upload your data (CSV, Excel, PDF or any file)")


uploaded_files = st.file_uploader(
    "",
    accept_multiple_files=True,
    type=[
        "csv",
        "xlsx",
        "xls",
        "pdf",
        "txt",
        "json",
        "docx"
    ]
)


if uploaded_files:

    for file in uploaded_files:

        key=f"{file.name}:{file.size}"

        if key not in st.session_state.file_keys:

            try:

                data=load_file(
                    file.name,
                    file.getvalue()
                )

                st.session_state.files.append(data)
                st.session_state.file_keys.add(key)

                st.success(
                    f"Loaded {file.name}"
                )

            except Exception as e:

                st.error(e)



if st.session_state.files:

    names=", ".join(
        f.filename
        for f in st.session_state.files
    )

    st.markdown(
        f"""
        <div class='loaded-box'>
        <b>Loaded data:</b> {names}
        </div>
        """,
        unsafe_allow_html=True
    )


st.divider()


st.markdown(
"<div class='chat-title'>🤖 AI Chatbot - Ask Questions About Your Data</div>",
unsafe_allow_html=True
)



if not st.session_state.files:

    st.markdown(
    """
    <div class='load-first'>
    📁 <b>Please load data first</b>
    - Use the upload option to add a file
    </div>
    """,
    unsafe_allow_html=True
    )


else:

    for index, msg in enumerate(st.session_state.history):

        with st.chat_message(msg["role"]):
            st.write(msg["content"])
            if msg.get("docx_bytes"):
                st.download_button(
                    "Download DOCX",
                    data=msg["docx_bytes"],
                    file_name=msg["docx_filename"],
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    key=f"docx_history_{index}"
                )
            if msg.get("pdf_bytes"):
                st.download_button(
                    "Download PDF",
                    data=msg["pdf_bytes"],
                    file_name=msg["pdf_filename"],
                    mime="application/pdf",
                    key=f"pdf_history_{index}"
                )
            if msg.get("pptx_bytes"):
                st.download_button(
                    "Download PPTX",
                    data=msg["pptx_bytes"],
                    file_name=msg["pptx_filename"],
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    key=f"pptx_history_{index}"
                )



    question = st.chat_input(
        "Ask a question about your uploaded data"
    )


    if question:

        with st.chat_message("user"):
            st.write(question)


        context = build_context(
            st.session_state.files,
            question
        )


        prompt = f"""
DOCUMENT CONTENT:

{context}

QUESTION:

{question}

Answer only from the document context above. If relevant details are present,
display the main information from the matching policy, subpolicy, procedure or
form instead of giving only a generic description.
"""


        recent_history=[
            {
                "role": msg["role"],
                "content": msg["content"]
            }
            for msg in st.session_state.history[-6:]
        ]

        messages=[
            {
                "role":"system",
                "content":SYSTEM_PROMPT
            },
            *recent_history,
            {
                "role":"user",
                "content":prompt
            }
        ]


        docx_bytes = None
        pdf_bytes = None
        pptx_bytes = None
        docx_filename = None
        pdf_filename = None
        pptx_filename = None

        with st.chat_message("assistant"):

            try:

                with st.spinner(
                    "Analyzing document..."
                ):

                    if is_document_generation_request(question):
                        generation_prompt = f"""
DOCUMENT CONTENT:

{context}

EMPLOYEE REQUEST:

{question}

Retrieve the relevant policy details from the document context, understand
eligibility, duration, approval workflow, required fields and key policy points,
then generate the requested output structure.
"""

                        if is_presentation_request(question):
                            generation_messages=[
                                {
                                    "role":"system",
                                    "content":PRESENTATION_GENERATION_PROMPT
                                },
                                {
                                    "role":"user",
                                    "content":generation_prompt
                                }
                            ]

                            raw_answer=chat_completion(
                                generation_messages,
                                api_key=llm_api_key.strip() or os.getenv("NVIDIA_API_KEY"),
                                model=st.session_state.llm_model,
                                temperature=0.1,
                                max_tokens=1600
                            )
                            presentation_data=normalize_presentation_data(
                                extract_json_object(raw_answer),
                                question
                            )
                            answer=render_presentation_preview(presentation_data)
                            base_filename=filename_from_title(presentation_data["title"])
                            docx_filename=f"{base_filename}.docx"
                            pdf_filename=f"{base_filename}.pdf"
                            pptx_filename=f"{base_filename}.pptx"
                            docx_bytes=create_docx_presentation(presentation_data)
                            pdf_bytes=create_pdf_presentation(presentation_data)
                            pptx_bytes=create_pptx_presentation(presentation_data)
                        else:
                            generation_messages=[
                                {
                                    "role":"system",
                                    "content":FORM_GENERATION_PROMPT
                                },
                                {
                                    "role":"user",
                                    "content":generation_prompt
                                }
                            ]

                            raw_answer=chat_completion(
                                generation_messages,
                                api_key=llm_api_key.strip() or os.getenv("NVIDIA_API_KEY"),
                                model=st.session_state.llm_model,
                                temperature=0.1,
                                max_tokens=1200
                            )
                            form_data=normalize_form_data(
                                extract_json_object(raw_answer),
                                question
                            )
                            answer=render_form_preview(form_data)
                            base_filename=filename_from_title(form_data["title"])
                            docx_filename=f"{base_filename}.docx"
                            pdf_filename=f"{base_filename}.pdf"
                            docx_bytes=create_docx_form(form_data)
                            pdf_bytes=create_pdf_form(form_data)
                    else:
                        answer=chat_completion(
                            messages,
                            api_key=llm_api_key.strip() or os.getenv("NVIDIA_API_KEY"),
                            model=st.session_state.llm_model
                        )


                st.write(answer)
                if docx_bytes:
                    st.download_button(
                        "Download DOCX",
                        data=docx_bytes,
                        file_name=docx_filename,
                        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                        key=f"docx_current_{docx_filename}"
                    )
                if pdf_bytes:
                    st.download_button(
                        "Download PDF",
                        data=pdf_bytes,
                        file_name=pdf_filename,
                        mime="application/pdf",
                        key=f"pdf_current_{pdf_filename}"
                    )
                if pptx_bytes:
                    st.download_button(
                        "Download PPTX",
                        data=pptx_bytes,
                        file_name=pptx_filename,
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        key=f"pptx_current_{pptx_filename}"
                    )


            except AIClientError as e:

                answer=str(e)
                st.error(answer)
                if "timed out" in answer.lower():
                    st.info("This usually clears up on its own — try sending the question again.")


            except Exception:

                full_trace = traceback.format_exc()
                logging.exception("Unhandled error while generating answer")
                answer = (
                    "Something went wrong while generating a response. "
                    "Please try again; if it keeps happening, try a shorter question "
                    "or a smaller document."
                )
                st.error(answer)
                with st.expander("Technical details"):
                    st.code(full_trace)



        st.session_state.history.append(
            {
                "role":"user",
                "content":question
            }
        )

        st.session_state.history.append(
            {
                "role":"assistant",
                "content":answer,
                "docx_bytes":docx_bytes,
                "docx_filename":docx_filename,
                "pdf_bytes":pdf_bytes,
                "pdf_filename":pdf_filename,
                "pptx_bytes":pptx_bytes,
                "pptx_filename":pptx_filename
            }
        )